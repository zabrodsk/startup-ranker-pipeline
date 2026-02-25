"""Extract text from PowerPoint files with slide numbers."""

from pathlib import Path


def extract_pptx(path: str | Path) -> list[dict]:
    """Extract text from a PPTX file, returning one item per slide.

    Returns:
        List of dicts with keys: text, slide (1-indexed), source_file.
    """
    from pptx import Presentation

    path = Path(path)
    prs = Presentation(str(path))
    results: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)

        combined = "\n".join(texts).strip()
        if combined:
            results.append({
                "text": combined,
                "slide": slide_num,
                "source_file": path.name,
            })

    return results
